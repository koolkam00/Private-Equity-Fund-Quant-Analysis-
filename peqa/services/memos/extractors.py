from __future__ import annotations

import base64
import logging
import os
import re
from io import BytesIO
from pathlib import Path

from flask import current_app

from peqa.services.memos.types import ExtractedDocument, ExtractedPage


logger = logging.getLogger(__name__)


def _native_pdf_empty_error() -> RuntimeError:
    return RuntimeError(
        "No extractable text was found in this PDF. If it is a scanned document, enable memo OCR or upload a text-based PDF."
    )


def _extract_pdf_pages_native(file_bytes: bytes) -> list[ExtractedPage]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf is required to extract PDF memo documents") from exc

    try:
        reader = PdfReader(BytesIO(file_bytes))
    except Exception as exc:
        raise RuntimeError(
            "Unable to read this PDF. If it is encrypted, password-protected, or uses unsupported security settings, "
            "export an unlocked PDF and upload it again."
        ) from exc

    try:
        if getattr(reader, "is_encrypted", False):
            try:
                decrypt_result = reader.decrypt("")
            except Exception as exc:
                raise RuntimeError(
                    "This PDF appears to be encrypted. Upload an unlocked PDF or a version that can be opened without a password."
                ) from exc
            if decrypt_result == 0:
                raise RuntimeError(
                    "This PDF is password-protected. Upload an unlocked PDF or a version that can be opened without a password."
                )
    except RuntimeError:
        raise
    except Exception as exc:
        raise RuntimeError(
            "Unable to read this PDF. If it is encrypted, password-protected, or uses unsupported security settings, "
            "export an unlocked PDF and upload it again."
        ) from exc

    pages = []
    for index, page in enumerate(reader.pages, start=1):
        pages.append(ExtractedPage(page_number=index, text=(page.extract_text() or "").strip()))
    return pages


def _ocr_enabled() -> bool:
    return bool(current_app.config.get("MEMO_ENABLE_OCR"))


def _page_text_length(text: str) -> int:
    return len(re.sub(r"\s+", "", text or ""))


def _pages_requiring_ocr(pages: list[ExtractedPage]) -> list[int]:
    minimum_chars = int(current_app.config.get("MEMO_OCR_MIN_PAGE_TEXT_CHARS", 80))
    return [
        page.page_number
        for page in pages
        if _page_text_length(page.text) < minimum_chars
    ]


def _render_pdf_pages_for_ocr(file_bytes: bytes, page_numbers: list[int]) -> dict[int, bytes]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required to OCR scanned PDF memo documents") from exc

    dpi = max(72, int(current_app.config.get("MEMO_OCR_RENDER_DPI", 180)))
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    rendered_pages: dict[int, bytes] = {}

    with fitz.open(stream=file_bytes, filetype="pdf") as pdf_document:
        for page_number in page_numbers:
            page = pdf_document.load_page(page_number - 1)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            rendered_pages[page_number] = pixmap.tobytes("png")
    return rendered_pages


def _ocr_page_image(page_number: int, png_bytes: bytes) -> str:
    try:
        from openai import OpenAI
    except ImportError as exc:
        raise RuntimeError("OpenAI SDK is required to OCR scanned PDF memo documents") from exc

    model = current_app.config.get("MEMO_OCR_MODEL") or "gpt-4.1-mini"
    image_b64 = base64.b64encode(png_bytes).decode("ascii")
    image_url = f"data:image/png;base64,{image_b64}"
    prompt = (
        "Transcribe all readable text from this investment document page. Preserve headings, bullets, and table values as plain text. "
        "Return only the transcription. If no readable text is present, return an empty string."
    )

    client = OpenAI()
    response = client.responses.create(
        model=model,
        input=[
            {
                "role": "user",
                "content": [
                    {"type": "input_text", "text": prompt},
                    {"type": "input_image", "image_url": image_url},
                ],
            }
        ],
    )
    output_text = (getattr(response, "output_text", None) or "").strip()
    if output_text:
        return output_text

    fragments: list[str] = []
    for item in getattr(response, "output", []) or []:
        for content in getattr(item, "content", []) or []:
            if getattr(content, "type", "") == "output_text" and getattr(content, "text", ""):
                fragments.append(content.text)
    text = "\n".join(fragment.strip() for fragment in fragments if fragment and fragment.strip()).strip()
    if text:
        return text
    logger.warning("OCR returned no text for page %s using model %s", page_number, model)
    return ""


def _ocr_pdf_pages(file_bytes: bytes, page_numbers: list[int]) -> dict[int, str]:
    if not page_numbers:
        return {}
    if not _ocr_enabled():
        return {}
    if not os.environ.get("OPENAI_API_KEY"):
        raise RuntimeError(
            "Memo OCR is enabled but OPENAI_API_KEY is not configured. Add the API key or upload a text-based PDF."
        )

    max_pages = max(1, int(current_app.config.get("MEMO_OCR_MAX_PAGES", 25)))
    if len(page_numbers) > max_pages:
        raise RuntimeError(
            f"OCR is required for {len(page_numbers)} page(s), which exceeds MEMO_OCR_MAX_PAGES={max_pages}. "
            "Split the PDF or increase the OCR page limit."
        )

    rendered_pages = _render_pdf_pages_for_ocr(file_bytes, page_numbers)
    return {page_number: _ocr_page_image(page_number, rendered_pages[page_number]) for page_number in page_numbers}


def _extract_text_pdf(file_bytes: bytes) -> tuple[list[ExtractedPage], dict[str, object]]:
    pages = _extract_pdf_pages_native(file_bytes)
    metadata: dict[str, object] = {
        "ocr_attempted": False,
        "ocr_page_numbers": [],
        "ocr_completed_pages": 0,
        "ocr_skipped_reason": None,
    }
    pages_to_ocr = _pages_requiring_ocr(pages)

    if not pages_to_ocr:
        return pages, metadata

    if not _ocr_enabled():
        if len(pages_to_ocr) == len(pages):
            raise _native_pdf_empty_error()
        return pages, metadata

    if not os.environ.get("OPENAI_API_KEY"):
        if len(pages_to_ocr) == len(pages):
            raise RuntimeError(
                "Memo OCR is enabled but OPENAI_API_KEY is not configured. Add the API key or upload a text-based PDF."
            )
        metadata["ocr_skipped_reason"] = "missing_openai_api_key"
        logger.warning("Skipping OCR fallback because OPENAI_API_KEY is missing.")
        return pages, metadata

    metadata["ocr_attempted"] = True
    metadata["ocr_page_numbers"] = pages_to_ocr
    ocr_results = _ocr_pdf_pages(file_bytes, pages_to_ocr)
    completed_pages = 0
    merged_pages = []
    for page in pages:
        replacement = (ocr_results.get(page.page_number) or "").strip()
        if replacement:
            merged_pages.append(ExtractedPage(page_number=page.page_number, text=replacement))
            completed_pages += 1
        else:
            merged_pages.append(page)
    metadata["ocr_completed_pages"] = completed_pages

    if not any(page.text.strip() for page in merged_pages):
        raise RuntimeError(
            "OCR could not extract readable text from this PDF. Upload a clearer scan or a text-based PDF."
        )
    return merged_pages, metadata


def _extract_text_docx(file_bytes: bytes) -> list[ExtractedPage]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required to extract DOCX memo documents") from exc

    document = Document(BytesIO(file_bytes))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
    return [ExtractedPage(page_number=1, text=text)]


def _extract_text_pptx(file_bytes: bytes) -> list[ExtractedPage]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to extract PPTX memo documents") from exc

    deck = Presentation(BytesIO(file_bytes))
    pages = []
    for index, slide in enumerate(deck.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        pages.append(ExtractedPage(page_number=index, text="\n".join(texts).strip()))
    return pages


def extract_document(document, storage) -> ExtractedDocument:
    file_bytes = storage.get(document.storage_key)
    suffix = Path(document.file_name).suffix.lower()
    metadata = {"suffix": suffix}
    if suffix in {".txt", ".md"}:
        text = file_bytes.decode("utf-8", errors="ignore")
        pages = [ExtractedPage(page_number=1, text=text)]
    elif suffix == ".pdf":
        pages, pdf_metadata = _extract_text_pdf(file_bytes)
        metadata.update(pdf_metadata)
        text = "\n\n".join(page.text for page in pages if page.text)
    elif suffix == ".docx":
        pages = _extract_text_docx(file_bytes)
        text = "\n\n".join(page.text for page in pages if page.text)
    elif suffix == ".pptx":
        pages = _extract_text_pptx(file_bytes)
        text = "\n\n".join(page.text for page in pages if page.text)
    else:
        text = file_bytes.decode("utf-8", errors="ignore")
        pages = [ExtractedPage(page_number=1, text=text)]

    return ExtractedDocument(
        document_id=document.id,
        file_name=document.file_name,
        mime_type=document.mime_type,
        text=text.strip(),
        pages=pages,
        metadata=metadata,
    )
